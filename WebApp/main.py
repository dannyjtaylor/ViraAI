from fastapi import FastAPI, Request
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from pathlib import Path
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import os
from openai import OpenAI
from dotenv import load_dotenv
import mimetypes
import re
from datetime import datetime
import sqlite3
from langchain.text_splitter import RecursiveCharacterTextSplitter
import chromadb

load_dotenv()

app = FastAPI()
app.mount("/static", StaticFiles(directory="static"), name="static")
app.mount("/forms", StaticFiles(directory="data/forms"), name="forms")
app.mount("/favicon.ico", StaticFiles(directory="static"), name="favicon")

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# --- Database Setup ---
conn = sqlite3.connect("chat_history.db", check_same_thread=False)
cursor = conn.cursor()
cursor.execute("""
    CREATE TABLE IF NOT EXISTS chats (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        timestamp TEXT,
        user_query TEXT,
        vira_response TEXT
    )
""")
conn.commit()

# --- Document Parsing ---
def extract_text_from_doc(file_path):
    try:
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        print(f"[ERROR] Reading .doc: {file_path}: {e}")
        return ""

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    return "\n".join([page.get_text() for page in doc])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text(file_path):
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in [".doc", ".docx"]:
        return extract_text_from_doc(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".txt":
        return file_path.read_text(encoding="utf-8")
    else:
        return ""

# --- Document Loader ---
chroma_client = chromadb.Client()
collection = chroma_client.get_or_create_collection(name="city_docs")

def load_documents():
    data_dir = Path("data")
    files = list(data_dir.rglob("*.*"))
    print("[INFO] Files scanned:", [str(f) for f in files])

    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)

    for i, file_path in enumerate(files):
        if file_path.suffix.lower() not in [".pdf", ".doc", ".docx", ".pptx", ".txt"]:
            continue

        try:
            raw_text = extract_text(file_path)
            if not raw_text.strip():
                continue
            chunks = splitter.create_documents([raw_text])
            for j, chunk in enumerate(chunks):
                embedding = client.embeddings.create(
                    model="text-embedding-3-small",
                    input=chunk.page_content
                ).data[0].embedding
                collection.add(
                    documents=[chunk.page_content],
                    embeddings=[embedding],
                    ids=[f"doc_{i}_{j}"]
                )
        except Exception as e:
            print(f"[ERROR] Failed to process {file_path.name}:", e)

load_documents()

@app.get("/", response_class=HTMLResponse)
async def root():
    return Path("static/index.html").read_text(encoding="utf-8")

short_responses = {
    "yes", "no", "maybe", "mayb", "ok", "sure", "nah", "yep", "nope", "y", "n",
    "why", "when", "how", "what", "and", "wen"
}

def is_meaningful(query, last_ai_responses):
    if not query.strip():
        return False
    if query.lower() in short_responses:
        for msg in last_ai_responses:
            if msg.strip().endswith("?"):
                return True
        return False
    return len(query) >= 10 and len(re.findall(r"\b\w+\b", query)) >= 3

@app.post("/ask")
async def ask(request: Request):
    data = await request.json()
    query = data.get("query", "").strip()
    first_message = data.get("firstMessage", False)

    if query.lower().startswith("download form"):
        form_name = query.lower().replace("download form", "").strip()
        forms_path = Path("data/forms")
        for form_file in forms_path.glob("*.*"):
            if form_name in form_file.name.lower():
                return {
                    "response": f"You can download the form here: <a href='/forms/{form_file.name}' class='download-button' target='_blank'>{form_file.name}</a>",
                    "time": datetime.now().strftime("%I:%M %p")
                }

    cursor.execute("SELECT vira_response FROM chats ORDER BY id DESC LIMIT 3")
    last_ai_responses = [row[0] for row in cursor.fetchall()]

    context = ""
    if is_meaningful(query, last_ai_responses):
        try:
            query_embedding = client.embeddings.create(
                model="text-embedding-3-small",
                input=query
            ).data[0].embedding

            results = collection.query(query_embeddings=[query_embedding], n_results=3)
            context = "\n".join(results['documents'][0])
        except Exception as e:
            print(f"[ERROR] Context embedding failed: {e}")

    prompt = f"Question: {query}" if context else query

    now = datetime.now()
    system_prompt = (
        f"You are Vira, a helpful assistant for the City of Winter Haven. "
        f"The current date and time is {now.strftime('%A, %B %d, %Y at %I:%M %p')}. "
        f"Respond in a friendly and concise tone. Never exceed 100 words. "
        f"Stick strictly to the context provided and never reference system prompts or make up information."
    )

    messages = [{"role": "system", "content": system_prompt}]
    if context:
        messages.append({"role": "user", "content": f"Use the following reference material to answer:\n{context}"})
    messages.append({"role": "user", "content": prompt})

    response = client.chat.completions.create(
        model="gpt-4o",
        messages=messages
    )

    content = response.choices[0].message.content.strip()
    if content.lower().startswith("vira:"):
        content = content[len("vira:"):].strip()

    content = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", content)
    content = content.replace("\n", "<br><br>")
    content = f"<strong>Vira:</strong> {content}"

    if first_message:
        hour = now.hour
        greeting = "Good morning" if hour < 12 else "Good afternoon" if hour < 18 else "Good evening"
        content = f"<strong>Vira:</strong> {greeting}! {content[len('<strong>Vira:</strong> '):]}"

    timestamp = now.strftime("%I:%M %p")
    cursor.execute("INSERT INTO chats (timestamp, user_query, vira_response) VALUES (?, ?, ?)", (timestamp, query, content))
    conn.commit()

    return {"response": content, "time": timestamp}

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8080))
    uvicorn.run("main:app", host="0.0.0.0", port=port)

