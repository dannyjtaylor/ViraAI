import os
import re
from datetime import datetime
from pathlib import Path
from fastapi import FastAPI, Request
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi import UploadFile, File
import requests
from fastapi.middleware.cors import CORSMiddleware
from openai import OpenAI
from dotenv import load_dotenv
import sqlite3
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
from langchain.text_splitter import RecursiveCharacterTextSplitter
import chromadb

# Load environment
load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

app = FastAPI()

# Mount folders
app.mount("/static", StaticFiles(directory="static"), name="static")
app.mount("/forms", StaticFiles(directory="data/forms"), name="forms")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# SQLite DB
conn = sqlite3.connect("chat_history.db", check_same_thread=False)
cursor = conn.cursor()
cursor.execute("""
    CREATE TABLE IF NOT EXISTS chats (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        timestamp TEXT,
        user_query TEXT,
        vira_response TEXT
    )
""")
conn.commit()

# File extractors
def extract_text(file_path):
    ext = file_path.suffix.lower()
    try:
        if ext == ".pdf":
            return "\n".join(page.get_text() for page in fitz.open(file_path))
        elif ext in [".doc", ".docx"]:
            return "\n".join(p.text for p in Document(file_path).paragraphs)
        elif ext == ".pptx":
            prs = Presentation(file_path)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif ext == ".txt":
            return file_path.read_text(encoding="utf-8")
    except Exception as e:
        print(f"[ERROR] Failed reading {file_path.name}: {e}")
    return ""

# Embedding
chroma_client = chromadb.Client()
collection = chroma_client.get_or_create_collection(name="city_docs")

def load_documents():
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    for i, file_path in enumerate(Path("data").rglob("*.*")):
        if file_path.suffix.lower() not in [".pdf", ".doc", ".docx", ".pptx", ".txt"]:
            continue
        try:
            text = extract_text(file_path)
            if text.strip():
                chunks = splitter.create_documents([text])
                for j, chunk in enumerate(chunks):
                    embedding = client.embeddings.create(
                        model="text-embedding-3-small",
                        input=chunk.page_content
                    ).data[0].embedding
                    collection.add(
                        documents=[chunk.page_content],
                        embeddings=[embedding],
                        ids=[f"{file_path.name}_{i}_{j}"]
                    )
        except Exception as e:
            print(f"[ERROR] Embedding failed for {file_path.name}: {e}")

load_documents()

# Serve home page
@app.get("/", response_class=HTMLResponse)
async def root():
    return Path("static/index.html").read_text(encoding="utf-8")

short_responses = {"yes", "no", "maybe", "ok", "sure", "nah", "yep", "nope", "y", "n", "why", "when", "how", "what", "and"}

def is_meaningful(query, recent_responses):
    if not query.strip() or query.lower() in short_responses:
        return any(msg.strip().endswith("?") for msg in recent_responses)
    return len(query) >= 10 and len(re.findall(r"\b\w+\b", query)) >= 3

@app.post("/ask")
async def ask(request: Request):
    data = await request.json()
    query = data.get("query", "").strip()
    first_message = data.get("firstMessage", False)

    # Form download handler
    if query.lower().startswith("download form"):
        name = query.lower().replace("download form", "").strip()
        for f in Path("data/forms").glob("*.*"):
            if name in f.name.lower():
                return {
                    "response": f"You can download the form here: <a href='/forms/{f.name}' class='download-button' target='_blank'>{f.name}</a>",
                    "time": datetime.now().strftime("%I:%M %p")
                }

    cursor.execute("SELECT vira_response FROM chats ORDER BY id DESC LIMIT 3")
    recent = [r[0] for r in cursor.fetchall()]

    context = ""
    if is_meaningful(query, recent):
        try:
            query_embedding = client.embeddings.create(
                model="text-embedding-3-small",
                input=query
            ).data[0].embedding
            results = collection.query(query_embeddings=[query_embedding], n_results=3)
            context = "\n".join(results["documents"][0])
        except Exception as e:
            print(f"[ERROR] Context retrieval failed: {e}")

    now = datetime.now()
    system_prompt = (
        f"You are Vira, a helpful assistant for the City of Winter Haven. "
        f"Today is {now.strftime('%A, %B %d, %Y')} and the current time is {now.strftime('%I:%M %p')}. "
        f"Keep answers concise and friendly. "
        f"Keep your response related only to information about Winter Haven and the documents you are provided. "
        f"Limit your response to 150 words or less. Remember previous conversations and use them to inform your responses and ask clarifying questions if needed."
    )


    messages = [{"role": "system", "content": system_prompt}]
    if context:
        messages.append({"role": "user", "content": f"Use this context:\n{context}"})
    messages.append({"role": "user", "content": query})

    ai_response = client.chat.completions.create(
        model="gpt-4o",
        messages=messages
    ).choices[0].message.content.strip()

    if ai_response.lower().startswith("vira:"):
        ai_response = ai_response[len("vira:"):].strip()

    ai_response = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", ai_response)
    ai_response = ai_response.replace("\n", "<br><br>")
    ai_response = f"<strong>Vira:</strong> {ai_response}"

    if first_message:
        greeting = "Good morning" if now.hour < 12 else "Good afternoon" if now.hour < 18 else "Good evening"
        ai_response = f"<strong>Vira:</strong> {greeting}! {ai_response[len('<strong>Vira:</strong> '):]}"

    timestamp = now.strftime("%I:%M %p")
    cursor.execute("INSERT INTO chats (timestamp, user_query, vira_response) VALUES (?, ?, ?)",
                   (timestamp, query, ai_response))
    conn.commit()

    return {"response": ai_response, "time": timestamp}

@app.post("/transcribe")
async def transcribe_audio(audio: UploadFile = File(...)):
    try:
        audio_data = await audio.read()
        response = requests.post(
            "https://api.openai.com/v1/audio/transcriptions",
            headers={
                "Authorization": f"Bearer {os.getenv('OPENAI_API_KEY')}"
            },
            files={
                "file": (audio.filename, audio_data, audio.content_type),
                "model": (None, "whisper-1")
            }
        )
        result = response.json()
        return {"transcript": result.get("text", "")}
    except Exception as e:
        return {"error": str(e)}

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8080))
    uvicorn.run("main:app", host="0.0.0.0", port=port)
